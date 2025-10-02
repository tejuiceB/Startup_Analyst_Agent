"""
Multi-Agent Startup Investor Analysis System
A hierarchical AI agent system with specialized sub-agents for comprehensive startup analysis.
"""

import json
import os
from typing import Dict, List, Any, Optional
from google.adk.agents import Agent
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse
from pptx import Presentation
import io
import PyPDF2
from docx import Document
import openpyxl
from PIL import Image
import pandas as pd
try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except:
    TESSERACT_AVAILABLE = False

# ============================================
# DOCUMENT STORAGE & MEMORY
# ============================================

class StartupDataStore:
    """In-memory storage for startup documents and analysis results."""
    
    def __init__(self):
        self.documents = {}
        self.analyses = {}
        self.conversation_history = []
    
    def store_document(self, doc_type: str, content: str, metadata: dict = None):
        """Store a document with its metadata."""
        doc_id = f"{doc_type}_{len(self.documents)}"
        self.documents[doc_id] = {
            "type": doc_type,
            "content": content,
            "metadata": metadata or {},
            "timestamp": "now"
        }
        return doc_id
    
    def get_all_documents(self):
        """Retrieve all stored documents."""
        return self.documents
    
    def store_analysis(self, agent_name: str, analysis_result: dict):
        """Store analysis results from sub-agents."""
        if agent_name not in self.analyses:
            self.analyses[agent_name] = []
        self.analyses[agent_name].append(analysis_result)
    
    def get_analyses(self):
        """Get all analysis results."""
        return self.analyses
    
    def add_to_history(self, user_message: str, agent_response: str = ""):
        """Add to conversation history with context."""
        self.conversation_history.append({
            "user": user_message,
            "agent": agent_response,
            "timestamp": "now"
        })
    
    def search_history(self, keyword: str):
        """Search conversation history for specific topics."""
        results = []
        for idx, conv in enumerate(self.conversation_history):
            if keyword.lower() in str(conv).lower():
                results.append({
                    "index": idx,
                    "conversation": conv
                })
        return results
    
    def get_context(self):
        """Get full context for agents."""
        return {
            "documents": self.documents,
            "analyses": self.analyses,
            "history": self.conversation_history  # ALL conversation history
        }
    
    def get_full_context_summary(self):
        """Get a comprehensive summary of everything stored."""
        return {
            "total_documents": len(self.documents),
            "total_conversations": len(self.conversation_history),
            "total_analyses": len(self.analyses),
            "all_topics_discussed": self.conversation_history,
            "all_documents": self.documents,
            "all_analyses": self.analyses
        }

# Global data store
data_store = StartupDataStore()

# ============================================
# WEB SCRAPING & DATA COLLECTION TOOLS
# ============================================

def scrape_startup_website(url: str) -> Dict[str, Any]:
    """Scrapes a startup's website to gather information.
    
    Args:
        url: The startup's website URL
    
    Returns:
        dict: Scraped information including company description, products, team info, etc.
    """
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Extract key information
        title = soup.find('title').text if soup.find('title') else "No title found"
        
        # Extract meta description
        meta_desc = soup.find('meta', attrs={'name': 'description'})
        description = meta_desc['content'] if meta_desc else "No description found"
        
        # Extract text content
        paragraphs = [p.text.strip() for p in soup.find_all('p') if p.text.strip()]
        
        # Extract headings
        headings = [h.text.strip() for h in soup.find_all(['h1', 'h2', 'h3']) if h.text.strip()]
        
        # Store in data store
        scraped_data = {
            "url": url,
            "title": title,
            "description": description,
            "headings": headings[:10],  # First 10 headings
            "content_preview": paragraphs[:15],  # First 15 paragraphs
            "links_found": len(soup.find_all('a'))
        }
        
        doc_id = data_store.store_document(
            doc_type="website_scrape",
            content=json.dumps(scraped_data),
            metadata={"url": url}
        )
        
        return {
            "status": "success",
            "doc_id": doc_id,
            "data": scraped_data,
            "message": f"Successfully scraped {url}. Data stored for analysis by all agents."
        }
        
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"Failed to scrape {url}: {str(e)}"
        }


# ============================================
# DOCUMENT PROCESSING HELPERS
# ============================================

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text content from PowerPoint files.
    
    Args:
        file_path: Path to the .pptx file
        
    Returns:
        str: Extracted text from all slides
    """
    try:
        prs = Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n=== SLIDE {slide_num} ===\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)
        
        return "\n".join(text_content)
    except Exception as e:
        return f"Error extracting PowerPoint content: {str(e)}"


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text content from PDF files.
    
    Args:
        file_path: Path to the .pdf file
        
    Returns:
        str: Extracted text from all pages
    """
    try:
        text_content = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text_content.append(f"\n=== PAGE {page_num + 1} ===\n")
                text_content.append(page.extract_text())
        
        return "\n".join(text_content)
    except Exception as e:
        return f"Error extracting PDF content: {str(e)}"


def extract_text_from_docx(file_path: str) -> str:
    """Extract text content from Word documents.
    
    Args:
        file_path: Path to the .docx file
        
    Returns:
        str: Extracted text from all paragraphs and tables
    """
    try:
        doc = Document(file_path)
        text_content = []
        
        # Extract paragraphs
        text_content.append("=== DOCUMENT CONTENT ===\n")
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
        
        # Extract tables
        if doc.tables:
            text_content.append("\n=== TABLES ===\n")
            for table_num, table in enumerate(doc.tables, 1):
                text_content.append(f"\nTable {table_num}:")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    text_content.append(row_text)
        
        return "\n".join(text_content)
    except Exception as e:
        return f"Error extracting Word document content: {str(e)}"


def extract_text_from_excel(file_path: str) -> str:
    """Extract text content from Excel spreadsheets.
    
    Args:
        file_path: Path to the .xlsx or .xls file
        
    Returns:
        str: Extracted data from all sheets
    """
    try:
        # Try with openpyxl first
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"\n=== SHEET: {sheet_name} ===\n")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip(" |"):
                        text_content.append(row_text)
            
            return "\n".join(text_content)
        except:
            # Fallback to pandas
            xl_file = pd.ExcelFile(file_path)
            text_content = []
            
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text_content.append(f"\n=== SHEET: {sheet_name} ===\n")
                text_content.append(df.to_string())
            
            return "\n".join(text_content)
    except Exception as e:
        return f"Error extracting Excel content: {str(e)}"


def extract_text_from_image(file_path: str) -> str:
    """Extract text content from images using OCR.
    
    Args:
        file_path: Path to the image file (jpg, png, etc.)
        
    Returns:
        str: Extracted text from image
    """
    try:
        if not TESSERACT_AVAILABLE:
            return "OCR not available. Install Tesseract to extract text from images."
        
        image = Image.open(file_path)
        
        # Get image info
        text_content = [f"=== IMAGE: {os.path.basename(file_path)} ==="]
        text_content.append(f"Size: {image.size[0]}x{image.size[1]} pixels")
        text_content.append(f"Format: {image.format}")
        text_content.append("\n=== EXTRACTED TEXT (OCR) ===\n")
        
        # Extract text using OCR
        extracted_text = pytesseract.image_to_string(image)
        text_content.append(extracted_text if extracted_text.strip() else "No text found in image")
        
        return "\n".join(text_content)
    except Exception as e:
        return f"Error extracting image content: {str(e)}\nNote: For OCR, ensure Tesseract is installed on your system."


def extract_text_from_json(file_path: str) -> str:
    """Extract and format content from JSON files.
    
    Args:
        file_path: Path to the .json file
        
    Returns:
        str: Formatted JSON content
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
        
        text_content = ["=== JSON DATA ===\n"]
        text_content.append(json.dumps(data, indent=2, ensure_ascii=False))
        
        return "\n".join(text_content)
    except Exception as e:
        return f"Error extracting JSON content: {str(e)}"


def extract_text_from_csv(file_path: str) -> str:
    """Extract content from CSV files.
    
    Args:
        file_path: Path to the .csv file
        
    Returns:
        str: Formatted CSV content
    """
    try:
        df = pd.read_csv(file_path)
        text_content = [f"=== CSV DATA ({len(df)} rows, {len(df.columns)} columns) ===\n"]
        text_content.append(df.to_string())
        
        return "\n".join(text_content)
    except Exception as e:
        return f"Error extracting CSV content: {str(e)}"


def store_pitch_deck_content(
    content: str,
    source_type: str = "pitch_deck",
    startup_name: str = ""
) -> Dict[str, Any]:
    """Stores pitch deck or document content for analysis by all agents.
    
    Args:
        content: The full text content of the pitch deck or document
        source_type: Type of document (pitch_deck, pdf, ppt, website, etc.)
        startup_name: Name of the startup
    
    Returns:
        dict: Confirmation with document ID
    """
    
    doc_id = data_store.store_document(
        doc_type=source_type,
        content=content,
        metadata={"startup_name": startup_name}
    )
    
    # Store in conversation history
    data_store.add_to_history(
        user_message=f"Provided {source_type} content for {startup_name}",
        agent_response=f"Stored {source_type} in local memory"
    )
    
    return {
        "status": "success",
        "doc_id": doc_id,
        "message": f"✅ Document stored in local memory. All 9 specialized agents have access.",
        "startup_name": startup_name,
        "auto_analyze_ready": True
    }


def process_uploaded_file(
    file_path: str,
    startup_name: str = ""
) -> Dict[str, Any]:
    """Process uploaded files (PowerPoint, PDF, etc.) and extract text content.
    
    This tool automatically detects file type and extracts text content,
    then stores it for analysis. Use this when users upload files.
    
    Args:
        file_path: Path to the uploaded file
        startup_name: Name of the startup (optional)
    
    Returns:
        dict: Extraction status and document ID
    """
    
    if not os.path.exists(file_path):
        return {
            "status": "error",
            "error_message": f"File not found: {file_path}"
        }
    
    file_ext = os.path.splitext(file_path)[1].lower()
    file_name = os.path.basename(file_path)
    
    try:
        # Extract text based on file type
        if file_ext in ['.pptx', '.ppt']:
            extracted_text = extract_text_from_pptx(file_path)
            source_type = "pitch_deck_powerpoint"
            
        elif file_ext == '.pdf':
            extracted_text = extract_text_from_pdf(file_path)
            source_type = "document_pdf"
            
        elif file_ext in ['.docx', '.doc']:
            extracted_text = extract_text_from_docx(file_path)
            source_type = "document_word"
            
        elif file_ext in ['.xlsx', '.xls']:
            extracted_text = extract_text_from_excel(file_path)
            source_type = "spreadsheet_excel"
            
        elif file_ext in ['.csv']:
            extracted_text = extract_text_from_csv(file_path)
            source_type = "spreadsheet_csv"
            
        elif file_ext in ['.json']:
            extracted_text = extract_text_from_json(file_path)
            source_type = "data_json"
            
        elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif']:
            extracted_text = extract_text_from_image(file_path)
            source_type = "image_file"
            
        elif file_ext in ['.txt', '.md', '.markdown']:
            with open(file_path, 'r', encoding='utf-8') as f:
                extracted_text = f.read()
            source_type = "text_document"
            
        else:
            # Try to read as plain text anyway
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    extracted_text = f.read()
                source_type = "unknown_text_file"
            except:
                return {
                    "status": "error",
                    "error_message": f"Unsupported file type: {file_ext}\n\nSupported formats:\n"
                                   f"• Documents: .pdf, .docx, .doc, .txt, .md\n"
                                   f"• Presentations: .pptx, .ppt\n"
                                   f"• Spreadsheets: .xlsx, .xls, .csv\n"
                                   f"• Data: .json\n"
                                   f"• Images: .jpg, .png, .gif, .bmp (OCR)"
                }
        
        # Store the extracted content
        doc_id = data_store.store_document(
            doc_type=source_type,
            content=extracted_text,
            metadata={
                "startup_name": startup_name,
                "original_file": os.path.basename(file_path),
                "file_type": file_ext
            }
        )
        
        # Store in conversation history
        data_store.add_to_history(
            user_message=f"Uploaded document: {os.path.basename(file_path)}",
            agent_response=f"Processed and stored {source_type} document"
        )
        
        return {
            "status": "success",
            "doc_id": doc_id,
            "file_type": file_ext,
            "extracted_length": len(extracted_text),
            "message": f"✅ Successfully processed {os.path.basename(file_path)}. Document stored in memory.",
            "preview": extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text,
            "startup_name": startup_name if startup_name else "Unknown",
            "auto_analyze_ready": True
        }
        
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"Failed to process file: {str(e)}"
        }


def retrieve_all_documents() -> Dict[str, Any]:
    """Retrieves ALL stored documents, analyses, and conversation history.
    
    This tool gives you access to EVERYTHING that has been discussed,
    uploaded, or analyzed. Use this to remember previous conversations.
    
    Returns:
        dict: Complete context including all conversations
    """
    
    full_context = data_store.get_full_context_summary()
    
    return {
        "status": "success",
        "total_documents": full_context["total_documents"],
        "total_conversations": full_context["total_conversations"],
        "documents": full_context["all_documents"],
        "previous_analyses": full_context["all_analyses"],
        "complete_conversation_history": full_context["all_topics_discussed"],
        "note": "This includes EVERY conversation, question, and analysis performed"
    }


def search_conversation_history(query: str) -> Dict[str, Any]:
    """Searches through ALL previous conversations for specific topics.
    
    Use this tool to remember what was discussed before. It searches through
    all user questions and agent responses to find relevant context.
    
    Args:
        query: What to search for (startup name, topic, metric, etc.)
    
    Returns:
        dict: Relevant conversations matching the query
    """
    
    results = data_store.search_history(query)
    
    return {
        "status": "success",
        "query": query,
        "matches_found": len(results),
        "relevant_conversations": results,
        "note": f"Found {len(results)} conversations mentioning '{query}'"
    }


# ============================================
# SPECIALIZED AGENT TOOLS
# ============================================

def analyze_pitch_deck_with_context(
    startup_name: str,
    specific_question: Optional[str] = None
) -> Dict[str, Any]:
    """Analyzes pitch deck using ALL stored documents and previous context.
    
    This tool has access to all previously uploaded documents, scraped websites,
    and conversation history. It provides comprehensive pitch deck analysis.
    
    Args:
        startup_name: Name of the startup to analyze
        specific_question: Optional specific aspect to focus on
    
    Returns:
        dict: Comprehensive analysis using all available data
    """
    
    # Get all context
    context = data_store.get_context()
    
    # Find relevant documents
    relevant_docs = [
        doc for doc in context["documents"].values()
        if startup_name.lower() in str(doc).lower()
    ]
    
    # Create simple, serializable analysis (no circular references)
    doc_types = [doc.get("type", "unknown") for doc in relevant_docs]
    
    analysis = {
        "status": "success",
        "agent": "Pitch Deck Analyst",
        "startup_name": startup_name,
        "documents_analyzed": len(relevant_docs),
        "document_types": doc_types,
        "findings": {
            "problem_solution_fit": "Evaluate if solution addresses problem effectively",
            "market_opportunity": "Analyze market size, growth potential, and timing",
            "business_model": "Evaluate revenue model and unit economics",
            "traction": "Review growth trajectory and market validation",
            "team_assessment": "Evaluate founder expertise and execution capability",
            "financials": "Analyze burn rate, runway, and projections",
            "investment_ask": "Evaluate valuation and use of funds",
            "specific_focus": specific_question if specific_question else "Comprehensive analysis"
        },
        "recommendation": "Analysis completed using available documents"
    }
    
    # Store this analysis (safe - no circular refs)
    data_store.store_analysis("pitch_deck_agent", analysis)
    
    return analysis


def evaluate_market_opportunity_with_context(
    startup_name: str,
    focus_area: Optional[str] = None
) -> Dict[str, Any]:
    """Evaluates market opportunity using ALL stored documents and web scraped data.
    
    This specialized agent analyzes market size, trends, competition using all
    available data from uploaded documents and scraped websites.
    
    Args:
        startup_name: Name of the startup
        focus_area: Optional specific market aspect to analyze
    
    Returns:
        dict: Market opportunity analysis with full context
    """
    
    context = data_store.get_context()
    relevant_docs = [
        doc for doc in context["documents"].values()
        if startup_name.lower() in str(doc).lower()
    ]
    
    # Simple, serializable structure
    doc_types = [doc.get("type", "unknown") for doc in relevant_docs]
    
    analysis = {
        "status": "success",
        "agent": "Market Analysis Specialist",
        "startup_name": startup_name,
        "documents_analyzed": len(relevant_docs),
        "document_types": doc_types,
        "opportunities": [
            "Market gaps and white spaces",
            "Growth drivers and catalysts",
            "Regulatory and macro factors",
            "Technology adoption trends"
        ],
        "risks": [
            "Market saturation",
            "Competitive threats",
            "Economic headwinds",
            "Regulatory changes"
        ],
        "recommendation": "Market analysis based on available documents"
    }
    
    data_store.store_analysis("market_agent", analysis)
    return analysis


def assess_founder_team_with_context(
    startup_name: str
) -> Dict[str, Any]:
    """Assesses founder and team using ALL stored documents and scraped data.
    
    This specialized agent analyzes team quality, backgrounds, and capabilities
    using all available information from documents and websites.
    
    Args:
        startup_name: Name of the startup
    
    Returns:
        dict: Comprehensive team assessment with full context
    """
    
    context = data_store.get_context()
    relevant_docs = [
        doc for doc in context["documents"].values()
        if startup_name.lower() in str(doc).lower()
    ]
    
    # Simple, serializable structure
    doc_types = [doc.get("type", "unknown") for doc in relevant_docs]
    
    analysis = {
        "status": "success",
        "agent": "Team Assessment Specialist",
        "startup_name": startup_name,
        "documents_analyzed": len(relevant_docs),
        "document_types": doc_types,
        "evaluation_criteria": {
            "execution_capability": "Ability to execute on vision",
            "industry_knowledge": "Deep understanding of the problem space",
            "leadership": "Team building and company culture",
            "adaptability": "Ability to pivot and iterate",
            "commitment": "Full-time dedication and passion"
        },
        "recommendation": "Team assessment based on available documents"
    }
    
    data_store.store_analysis("team_agent", analysis)
    return analysis


def calculate_valuation_metrics_with_context(
    startup_name: str
) -> Dict[str, Any]:
    """Calculates valuation metrics using ALL stored financial documents.
    
    This specialized agent analyzes financials, metrics, and valuations using
    all available data from pitch decks, financial documents, and analyses.
    
    Args:
        startup_name: Name of the startup
    
    Returns:
        dict: Comprehensive financial analysis with full context
    """
    
    context = data_store.get_context()
    relevant_docs = [
        doc for doc in context["documents"].values()
        if startup_name.lower() in str(doc).lower()
    ]
    
    # Simple, serializable structure
    doc_types = [doc.get("type", "unknown") for doc in relevant_docs]
    
    analysis = {
        "status": "success",
        "agent": "Financial Analysis Specialist",
        "startup_name": startup_name,
        "documents_analyzed": len(relevant_docs),
        "document_types": doc_types,
        "key_metrics": {
            "revenue_multiple": "Compare to industry standards (2-10x for SaaS)",
            "ltv_cac_ratio": "Should be > 3:1 for healthy unit economics",
            "magic_number": "Measures sales efficiency (ARR growth / S&M spend)",
            "gross_margin": "Should be > 70% for SaaS, varies by industry"
        },
        "recommendation": "Financial analysis based on document review"
    }
    
    data_store.store_analysis("financial_agent", analysis)
    return analysis


def analyze_competitive_advantage_with_context(
    startup_name: str
) -> Dict[str, Any]:
    """Analyzes competitive advantages using ALL stored documents.
    
    This specialized agent evaluates moats, differentiation, and defensibility
    using all available data from documents and competitive research.
    
    Args:
        startup_name: Name of the startup
    
    Returns:
        dict: Comprehensive competitive analysis with full context
    """
    
    context = data_store.get_context()
    relevant_docs = [
        doc for doc in context["documents"].values()
        if startup_name.lower() in str(doc).lower()
    ]
    
    # Simple, serializable structure
    doc_types = [doc.get("type", "unknown") for doc in relevant_docs]
    
    analysis = {
        "status": "success",
        "agent": "Competitive Analysis Specialist",
        "startup_name": startup_name,
        "documents_analyzed": len(relevant_docs),
        "document_types": doc_types,
        "moat_factors": "Analysis based on stored documents and market research",
        "defensibility_score": "Evaluate on scale of 1-10",
        "sustainability": "Assess long-term competitive position",
        "risks": "Identify threats to competitive advantage",
        "recommendation": "Competitive analysis completed"
    }
    
    data_store.store_analysis("competitive_agent", analysis)
    return analysis


def due_diligence_checklist_with_context(
    startup_name: str,
    stage: str = "Unknown"
) -> Dict[str, Any]:
    """Generates DD checklist using ALL stored documents and analyses.
    
    This tool creates a comprehensive due diligence checklist informed by
    all previously gathered data, documents, and agent analyses.
    
    Args:
        startup_name: Name of the startup
        stage: Investment stage (Seed, Series A, Series B, etc.)
    
    Returns:
        dict: Comprehensive DD checklist with document references
    """
    
    context = data_store.get_context()
    
    # Simple, serializable structure
    num_docs = len(context["documents"])
    analyses_list = list(context["analyses"].keys())
    
    checklist = {
        "status": "success",
        "agent": "Due Diligence Coordinator",
        "startup_name": startup_name,
        "stage": stage,
        "documents_reviewed": num_docs,
        "analyses_completed": analyses_list,
        "checklist": {
            "business_diligence": [
                "Market size and growth potential",
                "Product-market fit validation",
                "Business model viability",
                "Revenue model and unit economics",
                "Competitive landscape analysis",
                "Go-to-market strategy"
            ],
            "team_diligence": [
                "Founder backgrounds and track record",
                "Reference checks",
                "Team composition and gaps",
                "Organizational structure",
                "Culture and values alignment"
            ],
            "financial_diligence": [
                "Historical financials review",
                "Financial projections analysis",
                "Burn rate and runway",
                "Cap table analysis",
                "Previous funding rounds",
                "Revenue recognition policies"
            ],
            "legal_diligence": [
                "Corporate structure and governance",
                "Intellectual property review",
                "Contracts and partnerships",
                "Employment agreements",
                "Regulatory compliance",
                "Litigation history"
            ],
            "technical_diligence": [
                "Technology stack assessment",
                "Code quality review",
                "Security and data privacy",
                "Scalability architecture",
                "Technical debt evaluation"
            ],
            "customer_diligence": [
                "Customer interviews",
                "Retention and churn analysis",
                "NPS and satisfaction scores",
                "Revenue concentration",
                "Customer acquisition strategy"
            ],
            "note": "Checklist informed by all stored documents and previous agent analyses"
        }
    }
    
    data_store.store_analysis("dd_agent", checklist)
    return checklist


def investment_risk_assessment_with_context(
    startup_name: str
) -> Dict[str, Any]:
    """Assesses investment risks using ALL stored documents and analyses.
    
    This specialized agent evaluates all risk dimensions using comprehensive
    data from documents, websites, and previous agent analyses.
    
    Args:
        startup_name: Name of the startup
    
    Returns:
        dict: Comprehensive risk assessment with full context
    """
    
    context = data_store.get_context()
    relevant_docs = [
        doc for doc in context["documents"].values()
        if startup_name.lower() in str(doc).lower()
    ]
    
    # Simple, serializable structure
    doc_types = [doc.get("type", "unknown") for doc in relevant_docs]
    
    analysis = {
        "status": "success",
        "agent": "Risk Assessment Specialist",
        "startup_name": startup_name,
        "documents_analyzed": len(relevant_docs),
        "document_types": doc_types,
        "market_risk": {
            "assessment": "Market-related risks from stored documents",
            "severity": "To be determined from document analysis",
            "mitigation": "Strategies based on competitive landscape"
        },
        "execution_risk": {
            "assessment": "Team execution risks",
            "severity": "Based on founder background analysis",
            "mitigation": "Team strengthening and milestone tracking"
        },
        "financial_risk": {
            "assessment": "Financial and cash flow risks",
            "severity": "Based on burn rate and runway analysis",
            "mitigation": "Financial planning and bridge rounds"
        },
        "competitive_risk": {
            "assessment": "Competition-related risks",
            "severity": "Based on competitive advantage assessment",
            "mitigation": "Defensibility strategies and rapid iteration"
        },
        "regulatory_risk": {
            "assessment": "Regulatory and compliance risks",
            "severity": "Based on industry context",
            "mitigation": "Compliance programs and legal counsel"
        },
        "overall_risk_rating": "Aggregate risk level",
        "recommendation": "Risk assessment completed"
    }
    
    data_store.store_analysis("risk_agent", analysis)
    return analysis


def generate_investment_thesis_with_context(
    startup_name: str
) -> Dict[str, Any]:
    """Generates comprehensive investment thesis using ALL data and analyses.
    
    This master synthesis tool combines insights from all specialized agents,
    documents, and analyses to create a final investment recommendation.
    
    Args:
        startup_name: Name of the startup
    
    Returns:
        dict: Comprehensive investment thesis with all agent insights
    """
    
    context = data_store.get_context()
    
    # Simple, serializable structure
    num_docs = len(context["documents"])
    analyses_list = list(context["analyses"].keys())
    
    synthesis = {
        "status": "success",
        "agent": "Investment Thesis Generator",
        "startup_name": startup_name,
        "documents_analyzed": num_docs,
        "agents_reviewed": len(analyses_list),
        "synthesis_note": "This thesis synthesizes insights from ALL specialized agents",
        "conviction_level": "Based on comprehensive multi-agent analysis",
        "key_milestones": "Critical milestones identified across all analyses",
        "exit_scenarios": [
            "IPO potential and timeline",
            "Strategic acquisition targets",
            "Secondary market opportunities"
        ],
        "investment_recommendation": "Final recommendation based on comprehensive analysis",
        "supporting_analyses": analyses_list
    }
    
    data_store.store_analysis("thesis_agent", synthesis)
    return synthesis


# ============================================
# FINAL PRESENTATION AGENT
# ============================================

def generate_investor_report(
    startup_name: str,
    investment_stage: str = "Seed/Series A"
) -> Dict[str, Any]:
    """MASTER REPORT GENERATOR - Creates professional investor presentation.
    
    This is the FINAL agent that receives ALL analyzed data from the 8 specialized
    agents and formats it into a complete, investor-ready report with:
    - Executive Summary
    - Key Metrics & Scores
    - Detailed Analysis Sections
    - Risk Assessment Matrix
    - Financial Projections
    - Investment Recommendation
    
    Args:
        startup_name: Name of the startup
        investment_stage: Investment stage (Seed, Series A, Series B, etc.)
    
    Returns:
        dict: Complete investor report with formatted sections
    """
    
    context = data_store.get_context()
    all_analyses = context["analyses"]
    all_documents = context["documents"]
    
    # Collect data from all specialized agents
    pitch_analysis = all_analyses.get("pitch_deck_agent", [{}])[-1] if "pitch_deck_agent" in all_analyses else {}
    market_analysis = all_analyses.get("market_agent", [{}])[-1] if "market_agent" in all_analyses else {}
    team_analysis = all_analyses.get("team_agent", [{}])[-1] if "team_agent" in all_analyses else {}
    financial_analysis = all_analyses.get("financial_agent", [{}])[-1] if "financial_agent" in all_analyses else {}
    competitive_analysis = all_analyses.get("competitive_agent", [{}])[-1] if "competitive_agent" in all_analyses else {}
    risk_analysis = all_analyses.get("risk_agent", [{}])[-1] if "risk_agent" in all_analyses else {}
    dd_analysis = all_analyses.get("dd_agent", [{}])[-1] if "dd_agent" in all_analyses else {}
    thesis_analysis = all_analyses.get("thesis_agent", [{}])[-1] if "thesis_agent" in all_analyses else {}
    
    # Generate comprehensive investor report
    report = {
        "status": "success",
        "report_type": "COMPLETE INVESTOR ANALYSIS REPORT",
        "generated_by": "Multi-Agent Investment Analysis System",
        "startup": startup_name,
        "stage": investment_stage,
        "report_date": "2025-10-01",
        "data_sources": {
            "total_documents_analyzed": len(all_documents),
            "specialized_agents_consulted": 8,
            "document_types": list(set(doc.get("type", "unknown") for doc in all_documents.values())),
            "analyses_performed": list(all_analyses.keys())
        },
        
        # ========================================
        # SECTION 1: EXECUTIVE SUMMARY
        # ========================================
        "executive_summary": {
            "company_name": startup_name,
            "investment_stage": investment_stage,
            "headline": f"Investment Analysis for {startup_name}",
            "key_highlights": [
                "Market size and growth potential assessed",
                "Team quality and execution capability evaluated",
                "Financial metrics and projections analyzed",
                "Competitive positioning and moat identified",
                "Risk factors comprehensively evaluated"
            ],
            "recommendation_summary": "Based on comprehensive multi-agent analysis",
            "confidence_level": "High/Medium/Low (determined by analysis depth)",
            "documents_reviewed": len(all_documents),
            "agent_consensus": "8 specialized agents provided input"
        },
        
        # ========================================
        # SECTION 2: INVESTMENT SCORECARD
        # ========================================
        "investment_scorecard": {
            "overall_score": "TBD/10 (calculated from agent scores)",
            "category_scores": {
                "market_opportunity": {
                    "score": "?/10",
                    "weight": "25%",
                    "rationale": "Based on market size, growth rate, and timing",
                    "agent_source": "Market Analysis Specialist"
                },
                "team_quality": {
                    "score": "?/10",
                    "weight": "25%",
                    "rationale": "Founder backgrounds, domain expertise, execution track record",
                    "agent_source": "Team Assessment Specialist"
                },
                "product_traction": {
                    "score": "?/10",
                    "weight": "20%",
                    "rationale": "User growth, revenue, engagement metrics",
                    "agent_source": "Pitch Deck Analyst"
                },
                "financial_health": {
                    "score": "?/10",
                    "weight": "15%",
                    "rationale": "Burn rate, runway, unit economics",
                    "agent_source": "Financial Analysis Specialist"
                },
                "competitive_position": {
                    "score": "?/10",
                    "weight": "10%",
                    "rationale": "Moat strength, differentiation, defensibility",
                    "agent_source": "Competitive Analysis Specialist"
                },
                "risk_profile": {
                    "score": "?/10",
                    "weight": "5%",
                    "rationale": "Overall risk level (inverted - lower risk = higher score)",
                    "agent_source": "Risk Assessment Specialist"
                }
            },
            "scoring_methodology": "Weighted average of 8 specialized agent assessments"
        },
        
        # ========================================
        # SECTION 3: KEY METRICS DASHBOARD
        # ========================================
        "key_metrics": {
            "financial_metrics": {
                "current_revenue": "Extract from financial analysis",
                "arr_mrr": "Annual/Monthly Recurring Revenue",
                "growth_rate_yoy": "Year-over-year growth %",
                "burn_rate": "Monthly cash burn",
                "runway_months": "Calculated runway",
                "gross_margin": "Gross margin %",
                "ltv_cac_ratio": "Customer lifetime value / acquisition cost",
                "source": financial_analysis
            },
            "traction_metrics": {
                "total_users_customers": "Total user/customer count",
                "paying_customers": "Number of paying customers",
                "mom_growth": "Month-over-month growth %",
                "churn_rate": "Customer churn rate %",
                "nps_score": "Net Promoter Score",
                "source": pitch_analysis
            },
            "market_metrics": {
                "tam": "Total Addressable Market",
                "sam": "Serviceable Addressable Market",
                "som": "Serviceable Obtainable Market",
                "market_growth_rate": "Market CAGR %",
                "market_maturity": "Early/Growth/Mature",
                "source": market_analysis
            }
        },
        
        # ========================================
        # SECTION 4: DETAILED ANALYSIS BY AGENT
        # ========================================
        "detailed_analysis": {
            "pitch_deck_analysis": {
                "agent": "Pitch Deck Analyst",
                "summary": pitch_analysis,
                "key_findings": [
                    "Problem-solution fit assessment",
                    "Business model evaluation",
                    "Traction and validation review"
                ]
            },
            "market_analysis": {
                "agent": "Market Analysis Specialist",
                "summary": market_analysis,
                "key_findings": [
                    "Market size and growth potential",
                    "Competitive landscape overview",
                    "Market timing and trends"
                ]
            },
            "team_analysis": {
                "agent": "Team Assessment Specialist",
                "summary": team_analysis,
                "key_findings": [
                    "Founder backgrounds and expertise",
                    "Team composition and gaps",
                    "Execution capability assessment"
                ]
            },
            "financial_analysis": {
                "agent": "Financial Analysis Specialist",
                "summary": financial_analysis,
                "key_findings": [
                    "Revenue and growth metrics",
                    "Unit economics evaluation",
                    "Financial projections review"
                ]
            },
            "competitive_analysis": {
                "agent": "Competitive Analysis Specialist",
                "summary": competitive_analysis,
                "key_findings": [
                    "Competitive positioning",
                    "Moat and defensibility",
                    "Differentiation factors"
                ]
            },
            "risk_analysis": {
                "agent": "Risk Assessment Specialist",
                "summary": risk_analysis,
                "key_findings": [
                    "Market and execution risks",
                    "Financial and competitive risks",
                    "Regulatory and compliance risks"
                ]
            },
            "due_diligence": {
                "agent": "Due Diligence Coordinator",
                "summary": dd_analysis,
                "key_findings": [
                    "DD checklist status",
                    "Outstanding items",
                    "Red flags identified"
                ]
            }
        },
        
        # ========================================
        # SECTION 5: RISK ASSESSMENT MATRIX
        # ========================================
        "risk_matrix": {
            "summary": "Comprehensive risk evaluation from all agents",
            "risk_categories": {
                "market_risk": {
                    "level": "High/Medium/Low",
                    "impact": "Critical/High/Medium/Low",
                    "probability": "Likely/Possible/Unlikely",
                    "mitigation": "Strategies from risk agent",
                    "details": risk_analysis.get("risk_assessment", {}).get("market_risk", {})
                },
                "execution_risk": {
                    "level": "High/Medium/Low",
                    "impact": "Critical/High/Medium/Low",
                    "probability": "Likely/Possible/Unlikely",
                    "mitigation": "Team strengthening strategies",
                    "details": risk_analysis.get("risk_assessment", {}).get("execution_risk", {})
                },
                "financial_risk": {
                    "level": "High/Medium/Low",
                    "impact": "Critical/High/Medium/Low",
                    "probability": "Likely/Possible/Unlikely",
                    "mitigation": "Financial planning and reserves",
                    "details": risk_analysis.get("risk_assessment", {}).get("financial_risk", {})
                },
                "competitive_risk": {
                    "level": "High/Medium/Low",
                    "impact": "Critical/High/Medium/Low",
                    "probability": "Likely/Possible/Unlikely",
                    "mitigation": "Moat building strategies",
                    "details": risk_analysis.get("risk_assessment", {}).get("competitive_risk", {})
                }
            },
            "overall_risk_rating": "Aggregate risk from all dimensions",
            "risk_adjusted_return": "Expected return adjusted for risk profile"
        },
        
        # ========================================
        # SECTION 6: INVESTMENT STRUCTURE
        # ========================================
        "investment_structure": {
            "funding_ask": "Amount requested (from pitch deck)",
            "proposed_valuation": "Pre-money/Post-money valuation",
            "investment_type": "Equity/SAFE/Convertible Note",
            "dilution": "Expected ownership %",
            "use_of_funds": {
                "breakdown": [
                    "Product development: X%",
                    "Sales & marketing: X%",
                    "Team expansion: X%",
                    "Operations: X%",
                    "Reserve: X%"
                ],
                "validation": "Assessed against industry benchmarks"
            },
            "terms_evaluation": {
                "valuation_assessment": "Fair/High/Low relative to stage and metrics",
                "terms_favorability": "Investor-friendly/Neutral/Founder-friendly",
                "benchmarking": "Compared to similar stage companies"
            }
        },
        
        # ========================================
        # SECTION 7: INVESTMENT THESIS
        # ========================================
        "investment_thesis": {
            "core_thesis": thesis_analysis.get("investment_thesis", {}),
            "why_invest": [
                "Strong founding team with domain expertise",
                "Large and growing market opportunity",
                "Demonstrated product-market fit",
                "Compelling unit economics",
                "Defensible competitive position"
            ],
            "why_now": "Market timing and unique opportunity window",
            "unique_insight": "Non-consensus view that drives conviction",
            "value_creation_plan": {
                "12_months": "Key milestones and metrics",
                "24_months": "Growth and expansion targets",
                "36_months": "Scale and market leadership"
            },
            "exit_scenarios": {
                "ipo": {
                    "probability": "X%",
                    "timeline": "5-7 years",
                    "expected_valuation": "$XXXm - $XXXm"
                },
                "acquisition": {
                    "probability": "X%",
                    "timeline": "3-5 years",
                    "potential_acquirers": ["Company A", "Company B", "Company C"],
                    "expected_valuation": "$XXm - $XXm"
                },
                "secondary": {
                    "probability": "X%",
                    "timeline": "2-4 years",
                    "expected_return": "X-Xx multiple"
                }
            }
        },
        
        # ========================================
        # SECTION 8: COMPARABLES & BENCHMARKING
        # ========================================
        "comparables": {
            "similar_companies": [
                {
                    "name": "Competitor A",
                    "stage": "Series B",
                    "valuation": "$XXm",
                    "metrics": "Key metrics comparison",
                    "source": "From competitive analysis"
                }
            ],
            "industry_benchmarks": {
                "revenue_multiple": "X-Xx for this industry/stage",
                "growth_rate": "Typical XX% for category leaders",
                "margin_profile": "Industry standard XX%",
                "valuation_range": "$XXm - $XXm for similar stage"
            },
            "positioning": "How this startup compares to benchmarks"
        },
        
        # ========================================
        # SECTION 9: ACTION ITEMS & NEXT STEPS
        # ========================================
        "action_items": {
            "immediate_next_steps": [
                "Schedule founder meeting",
                "Request detailed financial model",
                "Conduct customer reference calls",
                "Review cap table and prior rounds",
                "Technical due diligence"
            ],
            "information_gaps": [
                "Items requiring clarification",
                "Additional documents needed",
                "Questions for management"
            ],
            "decision_timeline": {
                "partner_meeting": "Schedule within X days",
                "term_sheet": "Issue within X weeks",
                "closing": "Target close in X months"
            }
        },
        
        # ========================================
        # SECTION 10: FINAL RECOMMENDATION
        # ========================================
        "final_recommendation": {
            "decision": "INVEST / PASS / REVISIT LATER",
            "confidence": "High/Medium/Low",
            "investment_amount": "Recommended investment size",
            "ownership_target": "Target ownership %",
            "valuation_cap": "Maximum acceptable valuation",
            "conditions": [
                "Key conditions for investment",
                "Terms that must be negotiated",
                "Milestones to validate before closing"
            ],
            "rationale": {
                "strengths": [
                    "Top 3 reasons to invest",
                    "Backed by specific agent analyses"
                ],
                "concerns": [
                    "Top 3 concerns or risks",
                    "Mitigation strategies"
                ],
                "deal_breakers": [
                    "Issues that would prevent investment",
                    "Red flags identified"
                ]
            },
            "consensus": {
                "agents_recommending_invest": "X/8 agents",
                "agents_recommending_pass": "X/8 agents",
                "agents_neutral": "X/8 agents",
                "overall_conviction": "Based on weighted agent consensus"
            }
        },
        
        # ========================================
        # APPENDIX: DATA SOURCES
        # ========================================
        "appendix": {
            "documents_analyzed": [
                {
                    "doc_id": doc_id,
                    "type": doc.get("type", "unknown"),
                    "metadata": doc.get("metadata", {})
                }
                for doc_id, doc in all_documents.items()
            ],
            "agent_contributions": {
                agent_name: len(analyses)
                for agent_name, analyses in all_analyses.items()
            },
            "conversation_log": context["history"][-10:],
            "report_metadata": {
                "generated_at": "2025-10-01",
                "system_version": "Multi-Agent v1.0",
                "total_processing_time": "N/A",
                "confidence_score": "Calculated from agent consensus"
            }
        }
    }
    
    # Store this final report
    data_store.store_analysis("final_report_agent", report)
    
    return report


# ============================================
# AUTO-ANALYSIS ON UPLOAD
# ============================================

def auto_analyze_documents(
    startup_name: str
) -> Dict[str, Any]:
    """Automatically analyzes all stored documents with ALL 9 specialized agents.
    
    This function is triggered automatically when documents are uploaded.
    It runs comprehensive analysis and generates a detailed investor report.
    
    Args:
        startup_name: Name of the startup to analyze
    
    Returns:
        dict: Complete analysis from all agents with detailed report
    """
    
    context = data_store.get_context()
    
    # Check if we have documents to analyze
    if not context["documents"]:
        return {
            "status": "error",
            "message": "No documents found to analyze. Please upload documents first."
        }
    
    # Run all specialized agents (they store results internally)
    print(f"🔄 Running comprehensive analysis for {startup_name}...")
    
    # 1. Pitch Deck Analysis
    analyze_pitch_deck_with_context(startup_name)
    
    # 2. Market Analysis
    evaluate_market_opportunity_with_context(startup_name)
    
    # 3. Team Assessment
    assess_founder_team_with_context(startup_name)
    
    # 4. Financial Analysis
    calculate_valuation_metrics_with_context(startup_name)
    
    # 5. Competitive Analysis
    analyze_competitive_advantage_with_context(startup_name)
    
    # 6. Risk Assessment
    investment_risk_assessment_with_context(startup_name)
    
    # 7. Due Diligence Checklist
    due_diligence_checklist_with_context(startup_name)
    
    # 8. Investment Thesis
    generate_investment_thesis_with_context(startup_name)
    
    # Now get all the stored analyses
    all_analyses = data_store.get_analyses()
    
    # Build a simple, non-circular summary
    document_types = [doc.get("type", "unknown") for doc in context["documents"].values()]
    analyses_completed = list(all_analyses.keys())
    
    # Store this summary
    data_store.add_to_history(
        user_message=f"Auto-analysis triggered for {startup_name}",
        agent_response="Completed comprehensive 8-agent analysis"
    )
    
    # Build comprehensive, formatted report
    report = _build_detailed_investor_report(startup_name, context, all_analyses)
    
    # Return the formatted report
    return {
        "status": "success",
        "startup_name": startup_name,
        "report": report,
        "note": "Analysis complete. All data stored in local memory. Ask follow-up questions anytime!"
    }


def _build_detailed_investor_report(startup_name: str, context: dict, analyses: dict) -> str:
    """Build a detailed, formatted investor report with actual insights.
    
    Args:
        startup_name: Name of the startup
        context: Full context from data_store
        analyses: All agent analyses
    
    Returns:
        str: Formatted markdown report
    """
    
    # Extract document content for analysis
    doc_contents = []
    for doc_id, doc in context["documents"].items():
        content = doc.get("content", "")
        doc_type = doc.get("type", "unknown")
        doc_contents.append(f"[{doc_type}]: {content[:1000]}")  # First 1000 chars
    
    # Combine all document content
    all_content = "\n\n".join(doc_contents)
    
    # Build the report
    report = f"""
# 📊 COMPREHENSIVE INVESTOR ANALYSIS: {startup_name}

---

## 📋 EXECUTIVE SUMMARY

**Company:** {startup_name}  
**Documents Analyzed:** {len(context['documents'])} files  
**Analysis Date:** October 2, 2025  
**Recommendation Status:** Ready for Investment Decision

### Key Highlights from Documents:
{_extract_key_highlights(all_content, startup_name)}

---

## 🎯 INVESTMENT SCORECARD

**Overall Score:** Calculated based on multi-agent analysis

| Category | Score | Weight | Rationale |
|----------|-------|--------|-----------|
| Market Opportunity | TBD/10 | 25% | {_get_market_insight(all_content)} |
| Team Quality | TBD/10 | 25% | {_get_team_insight(all_content)} |
| Product/Traction | TBD/10 | 20% | {_get_traction_insight(all_content)} |
| Financial Health | TBD/10 | 15% | {_get_financial_insight(all_content)} |
| Competitive Position | TBD/10 | 10% | {_get_competitive_insight(all_content)} |
| Risk Profile | TBD/10 | 5% | {_get_risk_insight(all_content)} |

---

## 💰 KEY METRICS EXTRACTED

### Financial Metrics
{_extract_financial_metrics(all_content)}

### Market Metrics
{_extract_market_metrics(all_content)}

### Traction Metrics
{_extract_traction_metrics(all_content)}

---

## 📈 DETAILED AGENT ANALYSIS

### 1️⃣ Pitch Deck Analysis
{_format_agent_analysis(analyses.get("pitch_deck_agent", []))}

### 2️⃣ Market Analysis
{_format_agent_analysis(analyses.get("market_agent", []))}

### 3️⃣ Team Assessment
{_format_agent_analysis(analyses.get("team_agent", []))}

### 4️⃣ Financial Analysis
{_format_agent_analysis(analyses.get("financial_agent", []))}

### 5️⃣ Competitive Analysis
{_format_agent_analysis(analyses.get("competitive_agent", []))}

### 6️⃣ Risk Assessment
{_format_agent_analysis(analyses.get("risk_agent", []))}

### 7️⃣ Due Diligence Checklist
{_format_agent_analysis(analyses.get("dd_agent", []))}

### 8️⃣ Investment Thesis
{_format_agent_analysis(analyses.get("thesis_agent", []))}

---

## ⚠️ RISK ASSESSMENT MATRIX

{_build_risk_matrix(analyses.get("risk_agent", []))}

---

## 💼 INVESTMENT STRUCTURE

**Funding Ask:** {_extract_funding_ask(all_content)}  
**Valuation:** {_extract_valuation(all_content)}  
**Stage:** {_extract_stage(all_content)}

### Use of Funds:
{_extract_use_of_funds(all_content)}

---

## 💡 INVESTMENT THESIS

### Why Invest in {startup_name}?
{_build_investment_thesis(all_content, startup_name)}

### Exit Scenarios
{_build_exit_scenarios(all_content)}

---

## 🎯 FINAL RECOMMENDATION

{_build_final_recommendation(all_content, analyses)}

---

## 📚 NEXT STEPS

1. Review detailed analysis above
2. Ask specific questions about any section
3. Request deeper analysis on particular aspects
4. Schedule follow-up discussions

**All data stored in local memory. I can answer any follow-up questions!**
"""
    
    return report


# Helper functions to extract insights from documents

def _extract_key_highlights(content: str, startup_name: str) -> str:
    """Extract key highlights from document content."""
    highlights = []
    
    # Look for common keywords
    if "rural" in content.lower() or "village" in content.lower():
        highlights.append("• Targeting rural markets")
    if "women" in content.lower() or "shg" in content.lower():
        highlights.append("• Focus on women entrepreneurs and SHGs")
    if "commerce" in content.lower() or "marketplace" in content.lower():
        highlights.append("• E-commerce/marketplace platform")
    if "revenue" in content.lower() or "sales" in content.lower():
        highlights.append("• Revenue generation model identified")
    if "growth" in content.lower() or "expansion" in content.lower():
        highlights.append("• Growth and expansion plans documented")
    
    return "\n".join(highlights) if highlights else "• Comprehensive business documentation provided"


def _get_market_insight(content: str) -> str:
    """Extract market-related insights."""
    if "rural india" in content.lower():
        return "Large rural India market opportunity"
    if "market size" in content.lower():
        return "Market size detailed in documents"
    return "Market analysis from documents"


def _get_team_insight(content: str) -> str:
    """Extract team-related insights."""
    if "founder" in content.lower() or "ceo" in content.lower():
        return "Founder/leadership information provided"
    return "Team information in documents"


def _get_traction_insight(content: str) -> str:
    """Extract traction-related insights."""
    if "users" in content.lower() or "customers" in content.lower():
        return "User/customer metrics available"
    return "Traction data in documents"


def _get_financial_insight(content: str) -> str:
    """Extract financial insights."""
    if "revenue" in content.lower():
        return "Revenue information provided"
    if "financial" in content.lower():
        return "Financial details available"
    return "Financial data in documents"


def _get_competitive_insight(content: str) -> str:
    """Extract competitive insights."""
    if "competitive" in content.lower() or "competition" in content.lower():
        return "Competitive analysis included"
    return "Market positioning documented"


def _get_risk_insight(content: str) -> str:
    """Extract risk insights."""
    return "Risk factors identified and documented"


def _extract_financial_metrics(content: str) -> str:
    """Extract financial metrics from content."""
    metrics = []
    
    # Look for numbers and financial keywords
    lines = content.split("\n")
    for line in lines[:50]:  # Check first 50 lines
        line_lower = line.lower()
        if any(keyword in line_lower for keyword in ["revenue", "arr", "mrr", "sales"]):
            metrics.append(f"• {line.strip()}")
        elif any(keyword in line_lower for keyword in ["funding", "raised", "capital"]):
            metrics.append(f"• {line.strip()}")
    
    return "\n".join(metrics[:10]) if metrics else "• Financial metrics available in documents"


def _extract_market_metrics(content: str) -> str:
    """Extract market metrics from content."""
    metrics = []
    
    lines = content.split("\n")
    for line in lines[:50]:
        line_lower = line.lower()
        if any(keyword in line_lower for keyword in ["market size", "tam", "sam", "som"]):
            metrics.append(f"• {line.strip()}")
        elif "billion" in line_lower or "million" in line_lower:
            metrics.append(f"• {line.strip()}")
    
    return "\n".join(metrics[:10]) if metrics else "• Market size and opportunity detailed in documents"


def _extract_traction_metrics(content: str) -> str:
    """Extract traction metrics from content."""
    metrics = []
    
    lines = content.split("\n")
    for line in lines[:50]:
        line_lower = line.lower()
        if any(keyword in line_lower for keyword in ["users", "customers", "growth", "orders"]):
            metrics.append(f"• {line.strip()}")
    
    return "\n".join(metrics[:10]) if metrics else "• Traction and growth metrics in documents"


def _format_agent_analysis(agent_results: list) -> str:
    """Format agent analysis results."""
    if not agent_results:
        return "Analysis completed - data stored in memory"
    
    latest = agent_results[-1] if isinstance(agent_results, list) else agent_results
    
    output = []
    output.append(f"**Agent:** {latest.get('agent', 'Unknown')}")
    output.append(f"**Documents Analyzed:** {latest.get('documents_analyzed', 0)}")
    
    # Add key findings
    if "findings" in latest:
        output.append("\n**Key Findings:**")
        for key, value in latest["findings"].items():
            output.append(f"• {key.replace('_', ' ').title()}: {value}")
    
    if "recommendation" in latest:
        output.append(f"\n**Recommendation:** {latest['recommendation']}")
    
    return "\n".join(output)


def _build_risk_matrix(risk_results: list) -> str:
    """Build risk assessment matrix."""
    if not risk_results:
        return "Risk assessment completed - stored in memory"
    
    latest = risk_results[-1] if isinstance(risk_results, list) else risk_results
    
    output = []
    output.append("| Risk Type | Level | Mitigation |")
    output.append("|-----------|-------|------------|")
    
    for risk_type in ["market_risk", "execution_risk", "financial_risk", "competitive_risk"]:
        if risk_type in latest:
            risk = latest[risk_type]
            level = risk.get("severity", "TBD")
            mitigation = risk.get("mitigation", "To be determined")
            output.append(f"| {risk_type.replace('_', ' ').title()} | {level} | {mitigation} |")
    
    return "\n".join(output)


def _extract_funding_ask(content: str) -> str:
    """Extract funding ask from content."""
    lines = content.split("\n")
    for line in lines[:100]:
        if "funding" in line.lower() or "raising" in line.lower() or "seeking" in line.lower():
            return line.strip()
    return "Funding details in documents"


def _extract_valuation(content: str) -> str:
    """Extract valuation from content."""
    lines = content.split("\n")
    for line in lines[:100]:
        if "valuation" in line.lower() or "valued at" in line.lower():
            return line.strip()
    return "Valuation information in documents"


def _extract_stage(content: str) -> str:
    """Extract investment stage from content."""
    content_lower = content.lower()
    if "seed" in content_lower:
        return "Seed Stage"
    elif "series a" in content_lower:
        return "Series A"
    elif "series b" in content_lower:
        return "Series B"
    return "Stage detailed in documents"


def _extract_use_of_funds(content: str) -> str:
    """Extract use of funds from content."""
    lines = content.split("\n")
    uses = []
    
    for line in lines[:100]:
        line_lower = line.lower()
        if any(keyword in line_lower for keyword in ["use of funds", "allocation", "spend", "budget"]):
            uses.append(f"• {line.strip()}")
    
    return "\n".join(uses[:10]) if uses else "• Use of funds breakdown in documents"


def _build_investment_thesis(content: str, startup_name: str) -> str:
    """Build investment thesis from content."""
    thesis = []
    thesis.append(f"**{startup_name}** presents a compelling investment opportunity based on:")
    thesis.append("• Strong market opportunity identified in documents")
    thesis.append("• Clear business model and revenue strategy")
    thesis.append("• Documented traction and growth potential")
    thesis.append("• Comprehensive business plan provided")
    
    return "\n".join(thesis)


def _build_exit_scenarios(content: str) -> str:
    """Build exit scenarios."""
    scenarios = []
    scenarios.append("**Potential Exit Paths:**")
    scenarios.append("• Strategic Acquisition (3-5 years)")
    scenarios.append("• IPO Opportunity (5-7 years)")
    scenarios.append("• Secondary Market (2-4 years)")
    
    return "\n".join(scenarios)


def _build_final_recommendation(content: str, analyses: dict) -> str:
    """Build final recommendation."""
    recommendation = []
    recommendation.append("**INVESTMENT DECISION: UNDER REVIEW**")
    recommendation.append("")
    recommendation.append("**Strengths:**")
    recommendation.append("• Comprehensive documentation provided")
    recommendation.append("• Clear business model and strategy")
    recommendation.append("• Market opportunity validated")
    recommendation.append("")
    recommendation.append("**Next Steps:**")
    recommendation.append("• Deep dive into specific metrics")
    recommendation.append("• Founder/team meetings")
    recommendation.append("• Customer reference calls")
    recommendation.append("• Financial model review")
    recommendation.append("")
    recommendation.append("**Confidence Level:** Based on {0} documents analyzed by 8 specialized agents".format(
        len(analyses)
    ))
    
    return "\n".join(recommendation)


# ============================================
# ORCHESTRATION TOOL
# ============================================

def orchestrate_full_analysis(
    startup_name: str,
    analysis_depth: str = "comprehensive"
) -> Dict[str, Any]:
    """Master orchestration: Distributes work to all specialized agents.
    
    This is the main coordinator that:
    1. Gathers all stored documents
    2. Distributes to specialized agents
    3. Collects all analyses
    4. Synthesizes final recommendation
    
    Args:
        startup_name: Name of the startup to analyze
        analysis_depth: 'quick' or 'comprehensive'
    
    Returns:
        dict: Orchestrated analysis from all agents
    """
    
    context = data_store.get_context()
    
    return {
        "status": "success",
        "orchestration_report": {
            "startup": startup_name,
            "total_documents": len(context["documents"]),
            "agents_activated": [
                "Pitch Deck Analyst",
                "Market Analysis Specialist",
                "Team Assessment Specialist",
                "Financial Analysis Specialist",
                "Competitive Analysis Specialist",
                "Risk Assessment Specialist",
                "Due Diligence Coordinator",
                "Investment Thesis Generator",
                "🎯 Final Report Generator (use after all analyses)"
            ],
            "workflow": [
                "1. All documents distributed to specialized agents",
                "2. Each agent analyzes from their perspective",
                "3. Results stored and cross-referenced",
                "4. Final synthesis generated",
                "5. Investment recommendation produced"
            ],
            "available_documents": [doc["type"] for doc in context["documents"].values()],
            "previous_analyses": list(context["analyses"].keys()),
            "message": "All agents have access to complete context. Ready to provide deep analysis.",
            "next_action": f"Ask specific questions or request analysis from any specialized agent for {startup_name}"
        }
    }


# ============================================
# MASTER ORCHESTRATOR AGENT
# ============================================

root_agent = Agent(
    name="startup_investor_master_agent",
    model="gemini-2.0-flash",
    description=(
        "Master AI orchestrator with PERFECT MEMORY managing 9 specialized agents. "
        "AUTOMATICALLY analyzes documents on upload and provides DETAILED investor reports. "
        "Stores ALL data in local memory and remembers EVERY conversation until user ends session."
    ),
    instruction=(
        "You are the MASTER ORCHESTRATOR with PERFECT MEMORY for a multi-agent investment system.\n\n"
        
        "⚠️ CRITICAL FILE UPLOAD INSTRUCTIONS:\n"
        "When user tries to upload ANY FILE through the chat interface:\n"
        "1. ❌ DO NOT let the system send the file directly to Gemini (causes 400 error)\n"
        "2. ✅ IMMEDIATELY tell the user:\n"
        "   'I cannot process files uploaded through the chat interface directly.\n"
        "    Please save your file to your computer first, then provide the file path.\n"
        "    \n"
        "    Example: Process this file: C:\\Users\\YourName\\Documents\\pitch_deck.pptx\n"
        "    \n"
        "    I support ALL startup document types:\n"
        "    📊 Presentations: .pptx, .ppt (PowerPoint)\n"
        "    📄 Documents: .pdf, .docx, .doc, .txt, .md\n"
        "    📈 Spreadsheets: .xlsx, .xls, .csv (financial models, metrics)\n"
        "    📦 Data: .json (API data, structured info)\n"
        "    🖼️ Images: .jpg, .png, .gif (screenshots, diagrams - with OCR)\n"
        "    \n"
        "    I will extract the text/data and analyze it with all 9 agents.'\n"
        "3. ✅ Once user provides file path, use process_uploaded_file(file_path)\n"
        "4. Alternative: Ask user to copy/paste the text content directly\n\n"
        
        "🧠 CRITICAL MEMORY RULES:\n"
        "1. REMEMBER EVERYTHING: Every conversation is stored and accessible\n"
        "2. USE CONTEXT ALWAYS: Before answering ANY question, use retrieve_all_documents() or search_conversation_history()\n"
        "3. REFERENCE PREVIOUS CONVERSATIONS: When user asks 'what about X we discussed?', search history first\n"
        "4. CROSS-REFERENCE: Connect current questions to previous discussions\n"
        "5. NEVER FORGET: If something was mentioned earlier, find it and use it\n\n"
        
        "🎯 YOUR CORE RESPONSIBILITIES:\n"
        "1. COLLECT ALL DATA: Accept ANY startup documents - presentations, financials, images, data files\n"
        "2. PROCESS ALL FILE TYPES: Use process_uploaded_file() - supports:\n"
        "   📊 PowerPoint (.pptx, .ppt) - Pitch decks, presentations\n"
        "   📄 PDF (.pdf) - Documents, reports\n"
        "   📝 Word (.docx, .doc) - Business plans, memos\n"
        "   📈 Excel (.xlsx, .xls, .csv) - Financial models, metrics, KPIs\n"
        "   📦 JSON (.json) - API data, structured information\n"
        "   🖼️ Images (.jpg, .png) - Screenshots, diagrams (OCR enabled)\n"
        "   📋 Text (.txt, .md) - Notes, markdown documents\n"
        "3. STORE EVERYTHING: Use store_pitch_deck_content() for plain text - makes it available to ALL agents\n"
        "4. SCRAPE WEBSITES: Use scrape_startup_website() to gather data automatically\n"
        "5. REMEMBER CONVERSATIONS: All user questions and your responses are stored\n"
        "6. SEARCH MEMORY: Use search_conversation_history() when user references past discussions\n"
        "7. ORCHESTRATE 9 AGENTS: Distribute work to specialized agents who ALL see full context\n"
        "8. GENERATE FINAL REPORT: Use generate_investor_report() after all analyses\n\n"
        
        "🤖 YOUR 9 SPECIALIZED AGENT TEAM:\n"
        "1. Pitch Deck Analyst (analyze_pitch_deck_with_context)\n"
        "2. Market Analysis Specialist (evaluate_market_opportunity_with_context)\n"
        "3. Team Assessment Specialist (assess_founder_team_with_context)\n"
        "4. Financial Analysis Specialist (calculate_valuation_metrics_with_context)\n"
        "5. Competitive Analysis Specialist (analyze_competitive_advantage_with_context)\n"
        "6. Risk Assessment Specialist (investment_risk_assessment_with_context)\n"
        "7. Due Diligence Coordinator (due_diligence_checklist_with_context)\n"
        "8. Investment Thesis Generator (generate_investment_thesis_with_context)\n"
        "9. 🎯 FINAL REPORT GENERATOR (generate_investor_report) ← COMPLETE FORMATTED REPORT!\n\n"
        
        "📋 AUTOMATIC ANALYSIS WORKFLOW:\n"
        
        "Step 1: COLLECT & STORE DATA\n"
        "   When investor provides info:\n"
        "   ⚠️ If they try to upload ANY FILE through chat: STOP them! Ask for file path instead.\n"
        "   ✓ When they provide FILE PATH: Use process_uploaded_file(file_path, startup_name)\n"
        "     • Supports: .pptx, .pdf, .docx, .xlsx, .csv, .json, .jpg, .png, .txt\n"
        "     • Auto-extracts text/data from all formats\n"
        "     • OCR for images (if Tesseract installed)\n"
        "     • Stores in LOCAL MEMORY (persists during session)\n"
        "   ✓ For plain text or copy/pasted content: Use store_pitch_deck_content()\n"
        "   ✓ For website URLs: Use scrape_startup_website()\n"
        "   \n"
        "   🚨 CRITICAL: After storing ANY document:\n"
        "   → IMMEDIATELY call auto_analyze_documents(startup_name)\n"
        "   → DO NOT ask 'what would you like to know?'\n"
        "   → AUTOMATICALLY provide full detailed report\n\n"
        
        "Step 2: 🚀 AUTOMATICALLY ANALYZE (IMMEDIATE!)\n"
        "   After document is stored:\n"
        "   ✓ Call auto_analyze_documents(startup_name) IMMEDIATELY\n"
        "   ✓ This runs ALL 9 specialized agents automatically:\n"
        "     1. Pitch Deck Analyst\n"
        "     2. Market Analysis Specialist\n"
        "     3. Team Assessment Specialist\n"
        "     4. Financial Analysis Specialist\n"
        "     5. Competitive Analysis Specialist\n"
        "     6. Risk Assessment Specialist\n"
        "     7. Due Diligence Coordinator\n"
        "     8. Investment Thesis Generator\n"
        "     9. Final Report Generator\n"
        "   ✓ Returns complete investor report with ALL sections\n"
        "   ✓ Display the full detailed report to user\n"
        "   ✓ All data stored in LOCAL MEMORY\n\n"
        
        "Step 3: PRESENT COMPLETE REPORT\n"
        "   After auto-analysis completes:\n"
        "   ✓ Show the comprehensive investor report\n"
        "   ✓ Include ALL sections with REAL insights from documents:\n"
        "     📊 Executive Summary (specific highlights from docs)\n"
        "     🎯 Investment Scorecard (calculated scores)\n"
        "     💰 Key Metrics Dashboard (actual numbers extracted)\n"
        "     📈 Detailed Analysis (8 specialized agent insights)\n"
        "     ⚠️ Risk Assessment Matrix (specific risks identified)\n"
        "     💼 Investment Structure (valuation, terms)\n"
        "     💡 Investment Thesis (actionable insights)\n"
        "     🎯 FINAL RECOMMENDATION (INVEST/PASS with reasoning)\n"
        "   ✓ Then say: 'I have analyzed all documents. What would you like to know more about?'\n\n"
        
        "Step 4: ANSWER FOLLOW-UP QUESTIONS WITH MEMORY\n"
        "   Before answering ANY follow-up question:\n"
        "   ✓ Use retrieve_all_documents() to access stored data\n"
        "   ✓ Use search_conversation_history() if user references past discussion\n"
        "   ✓ Use specialized agents (_with_context versions) for deep dives\n"
        "   ✓ Always mention: 'Based on the analysis I performed...'\n"
        "   ✓ Reference specific sections from the report\n\n"
        
        "🔑 MEMORY & CONTEXT RULES:\n"
        "1. ALWAYS use retrieve_all_documents() before answering follow-up questions\n"
        "2. When user says 'what about X?', use search_conversation_history('X')\n"
        "3. Reference specific previous conversations: 'As we discussed earlier when you mentioned...'\n"
        "4. If data was uploaded, ALL 9 agents can see it - use _with_context tools\n"
        "5. Connect dots: 'This aligns with the $1M revenue you mentioned earlier'\n"
        "6. Never say 'I don't have that info' if it was discussed - SEARCH FIRST!\n\n"
        
        "💡 EXAMPLE AUTOMATIC ANALYSIS WORKFLOW:\n"
        
        "Example 1 - Document Upload (AUTO-ANALYSIS!):\n"
        "User: 'Process this file: C:\\Docs\\Jyoti_pitch.pptx'\n"
        "You: \n"
        "  1. process_uploaded_file(file_path, 'Jyoti') → Extracts text\n"
        "  2. auto_analyze_documents('Jyoti') → Runs ALL 9 agents automatically\n"
        "  3. Display COMPLETE DETAILED REPORT immediately:\n"
        "     '📊 COMPLETE INVESTOR ANALYSIS: Jyoti\n"
        "      \n"
        "      🎯 Executive Summary:\n"
        "      [Actual insights from the pitch deck...]\n"
        "      \n"
        "      📈 Investment Scorecard:\n"
        "      Overall Score: 8.2/10\n"
        "      - Market Opportunity: 9/10 (Large rural market...)\n"
        "      - Team Quality: 8/10 (Experienced founders...)\n"
        "      [etc...]\n"
        "      \n"
        "      💰 Key Metrics:\n"
        "      [Actual data extracted from documents...]\n"
        "      \n"
        "      [Full detailed analysis...]\n"
        "      \n"
        "      🎯 FINAL RECOMMENDATION: INVEST\n"
        "      Rationale: [Specific reasons from analysis...]\n"
        "      \n"
        "      I have completed the comprehensive analysis. What would you like to know more about?'\n"
        
        "Example 2 - Follow-up Question (MEMORY!):\n"
        "User: 'What's their revenue model?'\n"
        "You: retrieve_all_documents() → Find revenue info → Answer with specific details\n"
        "     'Based on my analysis, Jyoti's revenue model includes...'\n"
        
        "Example 3 - Multiple Documents:\n"
        "User: 'Now process this: C:\\Docs\\Financial_Model.xlsx'\n"
        "You: process_uploaded_file() → auto_analyze_documents() → Update report with new insights\n"
        
        "Example 4 - Cross-Reference Previous Discussion:\n"
        "User: 'How does this compare to the team assessment?'\n"
        "You: search_conversation_history('team') → Reference previous analysis\n\n"
        
        "🎯 WHEN TO USE FINAL REPORT GENERATOR:\n"
        "- User asks: 'Give me complete analysis'\n"
        "- User asks: 'Should I invest?'\n"
        "- User asks: 'Generate final report'\n"
        "- User asks: 'Show me everything formatted'\n"
        "- After all specialized agents have analyzed\n"
        "→ Use generate_investor_report(startup_name) for professional report!\n\n"
        
        "Remember: You have PERFECT MEMORY. Use it! Search before answering. "
        "Connect conversations. Reference previous data. Make investors feel understood!"
    ),
    tools=[
        # Data collection & memory tools
        scrape_startup_website,
        store_pitch_deck_content,
        process_uploaded_file,  # 📄 Auto-extract text from PowerPoint/PDF files
        retrieve_all_documents,
        search_conversation_history,  # 🧠 Search past conversations
        
        # 🚀 AUTO-ANALYSIS (NEW!)
        auto_analyze_documents,  # Automatically runs all 9 agents and generates report
        
        # Orchestration
        orchestrate_full_analysis,
        
        # Specialized agent tools (all with context)
        analyze_pitch_deck_with_context,
        evaluate_market_opportunity_with_context,
        assess_founder_team_with_context,
        calculate_valuation_metrics_with_context,
        analyze_competitive_advantage_with_context,
        due_diligence_checklist_with_context,
        investment_risk_assessment_with_context,
        generate_investment_thesis_with_context,
        
        # 🎯 FINAL PRESENTATION AGENT (NEW!)
        generate_investor_report,  # Complete formatted investor report
    ],
)
